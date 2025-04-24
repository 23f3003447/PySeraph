import speech_recognition as sr
import pyttsx3
import datetime
import webbrowser
import os
import shutil
import smtplib
import requests
import json
from email.message import EmailMessage
from pptx import Presentation
from fpdf import FPDF

# Initialize Text-to-Speech
engine = pyttsx3.init()

# DeepSeek API via OpenRouter
DEEPSEEK_API_KEY = "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"
DEEPSEEK_URL = "https://openrouter.ai/api/v1/chat/completions"

def speak(text):
    engine.say(text)
    engine.runAndWait()

def listen():
    recognizer = sr.Recognizer()
    with sr.Microphone() as source:
        print("🎤 Listening...")
        recognizer.adjust_for_ambient_noise(source)
        try:
            audio = recognizer.listen(source)
            command = recognizer.recognize_google(audio).lower()
            print(f"User said: {command}")
            return command
        except sr.UnknownValueError:
            speak("Sorry, I couldn't understand.")
        except sr.RequestError:
            speak("Error connecting to speech service.")
        return None

# DeepSeek AI Chat
def deepseek_chat(query):
    headers = {
        "Authorization": f"Bearer {DEEPSEEK_API_KEY}",
        "Content-Type": "application/json"
    }
    data = {
        "model": "deepseek/deepseek-chat",
        "messages": [{"role": "user", "content": query}]
    }
    try:
        response = requests.post(DEEPSEEK_URL, headers=headers, json=data)
        return response.json()["choices"][0]["message"]["content"]
    except Exception as e:
        return f"DeepSeek AI error: {e}"

def send_email(recipient, subject, body):
    sender_email = "pyseraphai@gmail.com"
    sender_password = "tpmp qotf wevy onih"
    msg = EmailMessage()
    msg['From'] = sender_email
    msg['To'] = recipient
    msg['Subject'] = subject
    msg.set_content(body)

    with smtplib.SMTP_SSL("smtp.gmail.com", 465) as server:
        server.login(sender_email, sender_password)
        server.send_message(msg)

    speak("Email has been sent successfully.")

def play_music():
    webbrowser.open("https://open.spotify.com")

def open_youtube():
    webbrowser.open("https://www.youtube.com")

def open_instagram():
    webbrowser.open("https://www.instagram.com")

def file_manager(command):
    if "new file" in command:
        file_name = input("Enter file name (with extension): ")
        with open(file_name, "w") as file:
            file.write("")
        speak(f"File {file_name} has been created.")
    elif "delete file" in command:
        file_name = input("Enter file name to delete: ")
        if os.path.exists(file_name):
            os.remove(file_name)
            speak(f"File {file_name} has been deleted.")
        else:
            speak("File does not exist.")
    elif "create folder" in command:
        folder_name = input("Enter folder name: ")
        os.makedirs(folder_name, exist_ok=True)
        speak(f"Folder {folder_name} has been created.")
    elif "delete folder" in command:
        folder_name = input("Enter folder name to delete: ")
        if os.path.exists(folder_name):
            shutil.rmtree(folder_name)
            speak(f"Folder {folder_name} has been deleted.")
        else:
            speak("Folder does not exist.")

def generate_ppt(topic):
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = topic
    content.text = "This is an auto-generated slide."
    prs.save(f"{topic}.pptx")
    speak(f"PowerPoint presentation '{topic}.pptx' created.")

def generate_pdf(content):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    pdf.multi_cell(200, 10, txt=content)
    pdf.output("generated_document.pdf")
    speak("PDF has been generated successfully.")

def get_weather(city="Mumbai"):
    api_key = "d26414e38d187576566cfbf22cd58b9e"
    url = f"https://api.openweathermap.org/data/2.5/weather?q={city}&appid={api_key}&units=metric"
    try:
        response = requests.get(url)
        data = response.json()
        print(json.dumps(data, indent=4))  # <-- JSON response shown here
        if response.status_code == 200 and data.get("cod") == 200:
            temp = data["main"]["temp"]
            desc = data["weather"][0]["description"]
            return f"The current temperature in {city} is {temp}°C with {desc}."
        else:
            return f"Error: {data.get('message', 'City not found')}"
    except Exception as e:
        return f"Weather error: {e}"

def get_news():
    api_key = "62ac0d8f472d433998615b2612a221cf"
    url = f"https://newsapi.org/v2/top-headlines?country=us&apiKey={api_key}"
    response = requests.get(url)

    if response.status_code == 200:
        news_data = response.json()
        articles = news_data.get("articles", [])[:5]
        if articles:
            speak("Here are the top news headlines:")
            for article in articles:
                print(article['title'])
                speak(article['title'])
        else:
            speak("No news articles found.")
    else:
        speak("Sorry, I couldn't fetch the news.")

def execute_command(command):
    if "time" in command:
        time_now = datetime.datetime.now().strftime("%I:%M %p")
        speak(f"The time is {time_now}")
    elif "search" in command:
        query = command.replace("search", "").strip()
        url = f"https://www.google.com/search?q={query}"
        webbrowser.open(url)
    elif "chat" in command:
        query = command.replace("chat", "").strip()
        response = deepseek_chat(query)
        speak(response)
    elif "email" in command:
        recipient = input("Enter recipient email: ")
        subject = input("Enter subject: ")
        body = input("Enter message: ")
        send_email(recipient, subject, body)
    elif "play music" in command:
        play_music()
    elif "open youtube" in command:
        open_youtube()
    elif "open instagram" in command:
        open_instagram()
    elif "file" in command:
        file_manager(command)
    elif "generate ppt" in command:
        topic = input("Enter topic: ")
        generate_ppt(topic)
    elif "generate pdf" in command:
        content = input("Enter content: ")
        generate_pdf(content)
    elif "weather" in command:
        city = input("Enter city: ")
        report = get_weather(city)
        speak(report)
    elif "news" in command:
        news = get_news()
        speak(news)
    elif "exit" in command or "quit" in command:
        speak("Goodbye! Have a great day.")
        exit()
    else:
        speak("Sorry, I didn't understand that command.")

if __name__ == "__main__":
    speak("Hello! PySeraph here. How can I help you?")
    while True:
        command = listen()
        if command:
            execute_command(command)